from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUTPUT_PATH = Path(__file__).resolve().parent / "product-ai-agent-client-presentation.pptx"

NAVY = RGBColor(23, 37, 68)
TEAL = RGBColor(44, 199, 195)
LIGHT = RGBColor(245, 248, 252)
SLATE = RGBColor(88, 102, 132)
DARK = RGBColor(31, 44, 68)
WHITE = RGBColor(255, 255, 255)
GREEN = RGBColor(57, 163, 114)
GOLD = RGBColor(214, 163, 52)


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    box = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(11.5), Inches(1.0))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = DARK
    if subtitle:
        p = tf.add_paragraph()
        p.text = subtitle
        p.font.size = Pt(11)
        p.font.color.rgb = SLATE


def add_footer(slide, text: str = "Product AI Agent") -> None:
    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(7.2), Inches(13.33), Inches(0.3))
    line.fill.solid()
    line.fill.fore_color.rgb = NAVY
    line.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.45), Inches(7.23), Inches(3.5), Inches(0.18))
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.color.rgb = WHITE


def add_bullets(slide, items: list[str], left: float, top: float, width: float, height: float, font_size: int = 18) -> None:
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK
        p.level = 0
        p.space_after = Pt(8)


def add_callout(slide, title: str, body: str, left: float, top: float, width: float, height: float, accent: RGBColor) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT
    shape.line.color.rgb = accent
    shape.line.width = Pt(1.6)
    tb = slide.shapes.add_textbox(Inches(left + 0.2), Inches(top + 0.18), Inches(width - 0.4), Inches(height - 0.36))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = accent
    p = tf.add_paragraph()
    p.text = body
    p.font.size = Pt(11)
    p.font.color.rgb = DARK


def add_process_row(slide, steps: list[tuple[str, str]]) -> None:
    left = 0.55
    width = 2.35
    top = 2.0
    height = 1.6
    gap = 0.18
    for index, (title, body) in enumerate(steps):
        x = left + index * (width + gap)
        card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(top), Inches(width), Inches(height))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT
        card.line.color.rgb = TEAL if index % 2 == 0 else NAVY
        card.line.width = Pt(1.5)
        num = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + 0.12), Inches(top + 0.12), Inches(0.38), Inches(0.38))
        num.fill.solid()
        num.fill.fore_color.rgb = TEAL if index % 2 == 0 else NAVY
        num.line.fill.background()
        tbn = slide.shapes.add_textbox(Inches(x + 0.21), Inches(top + 0.14), Inches(0.2), Inches(0.2))
        pn = tbn.text_frame.paragraphs[0]
        pn.text = str(index + 1)
        pn.font.size = Pt(11)
        pn.font.bold = True
        pn.font.color.rgb = WHITE
        pn.alignment = PP_ALIGN.CENTER
        tb = slide.shapes.add_textbox(Inches(x + 0.18), Inches(top + 0.55), Inches(width - 0.36), Inches(height - 0.68))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = DARK
        p = tf.add_paragraph()
        p.text = body
        p.font.size = Pt(10.5)
        p.font.color.rgb = SLATE


def build() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]

    # Slide 1
    slide = prs.slides.add_slide(blank)
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1.2))
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.7), Inches(1.45), Inches(8.8), Inches(1.4))
    p = tb.text_frame.paragraphs[0]
    p.text = "Product AI Agent"
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = DARK
    p = tb.text_frame.add_paragraph()
    p.text = "Market-informed product data generation, optimization, and compliant multi-marketplace image production"
    p.font.size = Pt(16)
    p.font.color.rgb = SLATE
    add_callout(slide, "Business Goal", "Turn a source product title and image into marketplace-ready data and visuals that are easier to publish, more competitive, and more consistent across channels.", 0.7, 3.2, 5.7, 2.05, TEAL)
    add_callout(slide, "Why This Matters", "Most tools only generate copy. This implementation combines product understanding, research, SEO, pricing guidance, image compliance, and second-pass optimization in one pipeline.", 6.7, 3.2, 5.7, 2.05, NAVY)
    add_footer(slide)

    # Slide 2
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Implementation Strategy", "Why this architecture was chosen instead of a single-prompt generator")
    add_bullets(
        slide,
        [
            "Separate product understanding from marketplace output generation",
            "Normalize attributes before writing marketplace content",
            "Use market research, SEO, and pricing as structured inputs",
            "Generate images through a cutout-first compliance workflow",
            "Run a second-pass optimizer to improve already generated data",
            "Validate outputs before product records are treated as final",
        ],
        0.9,
        1.5,
        5.8,
        4.8,
        17,
    )
    add_callout(slide, "Why not one prompt?", "Single-shot generators are fast but inconsistent. This layered design improves accuracy, repeatability, and control across multiple marketplaces.", 7.1, 1.9, 5.4, 1.55, GOLD)
    add_callout(slide, "Result", "A production-style workflow that is easier to extend with live marketplace integrations, validation rules, and publishing actions.", 7.1, 3.75, 5.4, 1.55, GREEN)
    add_footer(slide)

    # Slide 3
    slide = prs.slides.add_slide(blank)
    add_title(slide, "End-to-End Process", "How a product moves through the system")
    add_process_row(
        slide,
        [
            ("Input", "Source product title and product image are submitted."),
            ("Vision", "Image cues such as palette, brightness, and visible traits are analyzed."),
            ("Core Data", "Canonical product type, category, summary, features, and attributes are created."),
            ("Research", "Marketplace signals, comparable listings, and pricing ranges are built."),
            ("Output", "Marketplace data, images, validation, and stored product record are returned."),
        ],
    )
    add_footer(slide)

    # Slide 4
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Product Intelligence Layers", "The system uses multiple optimization layers rather than relying on raw text generation")
    add_callout(slide, "1. Vision Understanding", "Extracts product cues from the uploaded image to support category and attribute inference.", 0.8, 1.5, 3.9, 1.45, TEAL)
    add_callout(slide, "2. Attribute Mapping", "Normalizes raw attributes into cleaner marketplace-ready fields such as color, material, size, and style.", 4.75, 1.5, 3.9, 1.45, NAVY)
    add_callout(slide, "3. Market Research", "Builds marketplace-specific keyword signals, comparable titles, and pricing bands.", 8.7, 1.5, 3.9, 1.45, GOLD)
    add_callout(slide, "4. SEO Layer", "Generates primary keywords, secondary keywords, title terms, and marketplace keyword clusters.", 0.8, 3.35, 3.9, 1.45, NAVY)
    add_callout(slide, "5. Pricing Layer", "Creates recommended price, floor, ceiling, and strategy per marketplace.", 4.75, 3.35, 3.9, 1.45, GREEN)
    add_callout(slide, "6. Validation Layer", "Checks field quality, marketplace limits, and image compliance issues before final use.", 8.7, 3.35, 3.9, 1.45, TEAL)
    add_footer(slide)

    # Slide 5
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Marketplace-Specific Data Generation", "Why separate marketplace agents improve competitiveness")
    add_bullets(
        slide,
        [
            "Amazon: structured bullets, search terms, and clean white-background asset",
            "eBay: concise titles, item specifics, and comparable listing style",
            "Etsy: descriptive tags, materials, occasion, and gift-oriented search language",
            "TikTok Shop: short-form social commerce tone and vertical hero image output",
            "Shopify: storefront-ready merchandising copy, tags, and SEO metadata",
        ],
        0.8,
        1.45,
        6.0,
        4.7,
        16,
    )
    add_callout(slide, "Why it performs better", "Each marketplace rewards different behaviors. Generating dedicated outputs per channel improves listing fit, discoverability, and publishing readiness.", 7.1, 2.0, 5.35, 2.25, NAVY)
    add_footer(slide)

    # Slide 6
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Image Optimization Workflow", "Built to preserve product identity while adapting to marketplace image requirements")
    add_process_row(
        slide,
        [
            ("Source Save", "Original upload is retained for audit and regeneration."),
            ("Cutout", "Transparent product cutout is created first."),
            ("Amazon/eBay", "White-background compositions are built for compliance."),
            ("Shopify/TikTok/Etsy", "Styled backgrounds are added without changing product color or shape."),
            ("Validation", "Image dimensions, background expectations, and format issues are checked."),
        ],
    )
    add_footer(slide)

    # Slide 7
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Second-Pass Optimization API", "Why an optimizer endpoint was added after generation")
    add_bullets(
        slide,
        [
            "Generation creates the first complete product record",
            "Optimization reuses the stored record and recalculates research, SEO, and pricing context",
            "The optimizer improves marketplace titles, descriptions, keywords, attributes, and pricing language",
            "Images remain unchanged during optimization so product media stays stable",
            "Optimization can target all marketplaces or only selected marketplaces",
        ],
        0.85,
        1.55,
        6.1,
        4.8,
        16,
    )
    add_callout(slide, "Business Value", "This creates a workflow where teams can generate quickly, then refine strategically before publishing. It is closer to how real catalog operations work.", 7.15, 2.05, 5.3, 2.0, GREEN)
    add_footer(slide)

    # Slide 8
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Why This Is More Competitive Than Basic Tools", "Positioning against common market alternatives")
    add_bullets(
        slide,
        [
            "Basic tools usually stop at title + description generation",
            "This implementation adds structured attributes, SEO signals, pricing support, image compliance, and validation",
            "It supports multiple marketplaces through dedicated output logic",
            "It is designed for live marketplace adapters instead of hardcoding one data source",
            "It supports second-pass optimization, which most lightweight generators do not provide",
        ],
        0.85,
        1.55,
        6.0,
        4.8,
        16,
    )
    add_callout(slide, "Competitive Position", "The strength of this system is not only copy quality. It is the combination of research, structure, adaptability, optimization, and operational readiness.", 7.15, 2.0, 5.3, 2.05, GOLD)
    add_footer(slide)

    # Slide 9
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Current Scope and Upgrade Path", "Clear distinction between what is live now and what is ready for expansion")
    add_callout(slide, "Live Today", "End-to-end generation pipeline, second-pass optimization API, multi-marketplace data generation, compliant image workflow, and live official eBay research integration.", 0.8, 1.55, 5.9, 2.05, GREEN)
    add_callout(slide, "Ready to Extend", "Amazon, Etsy, TikTok Shop, and Shopify can be upgraded to live research mode with official API credentials and platform access.", 6.9, 1.55, 5.6, 2.05, NAVY)
    add_callout(slide, "Why this matters", "The architecture already supports the right direction: live marketplace adapters can be added without rebuilding the core product intelligence pipeline.", 0.8, 4.0, 11.7, 1.5, TEAL)
    add_footer(slide)

    prs.save(OUTPUT_PATH)


if __name__ == "__main__":
    build()
